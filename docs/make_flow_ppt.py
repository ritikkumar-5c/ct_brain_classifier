#!/usr/bin/env python3
"""
Generate a PowerPoint of the CT-brain MIL classifier slice-handling flows:
detailed pipeline diagrams for the CAPPED (default) and ALL-SLICES cases,
plus a side-by-side comparison. Output: ct_brain_flow.pptx
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from lxml import etree

A = "http://schemas.openxmlformats.org/drawingml/2006/main"

# ---- palette ----
BG      = RGBColor(0xF7, 0xF9, 0xFC)
TITLE   = RGBColor(0x10, 0x2A, 0x43)
COMMON  = RGBColor(0x2E, 0x5E, 0x8C)   # shared stages (blue)
DIFF    = RGBColor(0xC1, 0x57, 0x1E)   # stages that differ (orange)
IO      = RGBColor(0x3B, 0x7A, 0x57)   # data in/out (green)
WHITE   = RGBColor(0xFF, 0xFF, 0xFF)
GREY    = RGBColor(0x55, 0x55, 0x55)
NOTE    = RGBColor(0x33, 0x3A, 0x44)

KIND = {"common": COMMON, "diff": DIFF, "io": IO}

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]


def bg(slide):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = BG


def textbox(slide, x, y, w, h, text, size, color, bold=False, align=PP_ALIGN.LEFT,
            anchor=MSO_ANCHOR.TOP, italic=False):
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    for i, line in enumerate(text.split("\n")):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        r = p.add_run(); r.text = line
        r.font.size = Pt(size); r.font.bold = bold; r.font.italic = italic
        r.font.color.rgb = color
        r.font.name = "Calibri"
    return tb


def box(slide, x, y, w, h, title, kind, sub=None):
    shp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                 Inches(x), Inches(y), Inches(w), Inches(h))
    shp.fill.solid(); shp.fill.fore_color.rgb = KIND[kind]
    shp.line.color.rgb = WHITE; shp.line.width = Pt(1)
    shp.shadow.inherit = False
    tf = shp.text_frame; tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.margin_top = Pt(2); tf.margin_bottom = Pt(2)
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = title
    r.font.size = Pt(12.5); r.font.bold = True; r.font.color.rgb = WHITE
    r.font.name = "Calibri"
    if sub:
        p2 = tf.add_paragraph(); p2.alignment = PP_ALIGN.CENTER
        r2 = p2.add_run(); r2.text = sub
        r2.font.size = Pt(9.5); r2.font.color.rgb = RGBColor(0xE9, 0xEE, 0xF5)
        r2.font.name = "Calibri"
    return shp


def arrow(slide, x1, y1, x2, y2):
    c = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT,
                                   Inches(x1), Inches(y1), Inches(x2), Inches(y2))
    c.line.color.rgb = GREY; c.line.width = Pt(2.25)
    ln = c.line._get_or_add_ln()
    tail = etree.SubElement(ln, f"{{{A}}}tailEnd")
    tail.set("type", "triangle"); tail.set("w", "med"); tail.set("len", "med")
    return c


def legend(slide, x, y):
    items = [("Shared stage", COMMON), ("Differs by mode", DIFF), ("Data in / out", IO)]
    for i, (lab, col) in enumerate(items):
        sw = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x + i * 2.05), Inches(y),
                                    Inches(0.16), Inches(0.16))
        sw.fill.solid(); sw.fill.fore_color.rgb = col; sw.line.fill.background()
        sw.shadow.inherit = False
        textbox(slide, x + 0.22 + i * 2.05, y - 0.07, 1.8, 0.3, lab, 10, NOTE)


def pipeline_slide(title, subtitle, stages):
    """stages: list of (title, sub, kind, note)."""
    s = prs.slides.add_slide(BLANK); bg(s)
    textbox(s, 0.5, 0.18, 12.4, 0.5, title, 26, TITLE, bold=True)
    textbox(s, 0.5, 0.74, 12.4, 0.35, subtitle, 13, NOTE, italic=True)
    legend(s, 9.0, 0.34)

    bx, bw = 0.85, 5.1
    top, bh, gap = 1.28, 0.50, 0.135
    nx = bx + bw + 0.45
    centers = []
    for i, (t, sub, kind, note) in enumerate(stages):
        y = top + i * (bh + gap)
        box(s, bx, y, bw, bh, t, kind, sub)
        centers.append((bx + bw / 2, y, y + bh))
        if note:
            textbox(s, nx, y - 0.02, 6.7, bh + 0.1, note, 10.5, NOTE,
                    anchor=MSO_ANCHOR.MIDDLE)
    for i in range(len(stages) - 1):
        cx, _, ybot = centers[i]
        _, ytop, _ = centers[i + 1]
        arrow(s, cx, ybot, cx, ytop)
    return s


# ---------------- Slide 1: title ----------------
s = prs.slides.add_slide(BLANK); bg(s)
textbox(s, 0.8, 2.3, 11.7, 1.2,
        "CT-Brain MIL Classifier", 40, TITLE, bold=True, align=PP_ALIGN.CENTER)
textbox(s, 0.8, 3.4, 11.7, 0.8,
        "Slice-handling flow — Capped (default) vs All-slices",
        22, COMMON, align=PP_ALIGN.CENTER)
textbox(s, 0.8, 4.5, 11.7, 1.2,
        "MaxViT backbone  +  gated-attention MIL pooling  •  one label per studyID_seriesID set\n"
        "3-channel input: brain (W80/L40) · subdural (W200/L80) · bone (W2800/L600)  →  "
        "normal / near-normal / abnormal",
        14, NOTE, align=PP_ALIGN.CENTER)

# ---------------- Slide 2: capped ----------------
capped = [
    ("DICOM set", "prepared_brain/<studyID>_<seriesID>/", "io",
     "N slices per set — VARIES (18 … 300+). One label per set."),
    ("index_studies()", "(study_id, [slice paths], label)", "common",
     "labels.csv: study_id,label  (normal / near-normal / abnormal)"),
    ("Slice selection — CAPPED", "TRAIN: random 24   •   EVAL: even-spaced 32", "diff",
     "K = min(cap, N).  Random subsample = augmentation; even spacing = deterministic eval."),
    ("Per-slice render → 3 channels", "HU → brain · subdural · bone windows", "common",
     "resize 224×224 + ImageNet norm  →  one slice = 3 × 224 × 224"),
    ("Bag = stack(K slices)", "K × 3 × 224 × 224  (K ≤ 32)", "common",
     "Anatomical order (sorted by InstanceNumber)."),
    ("mil_collate — pad + mask", "B × maxK × 3 × 224 × 224  +  mask B×maxK", "common",
     "Ragged bags padded to batch max; mask marks real vs padded slices."),
    ("MaxViT encoder — single forward", "backbone(B·K images) → B × K × D", "diff",
     "ALL slices of the batch encoded in ONE pass (memory ∝ B·K → why K is capped)."),
    ("Gated-attention pooling", "mask→ -inf, softmax, weighted sum → B × D", "common",
     "Collapses K slices → 1 study vector. Padding gets 0 weight. Yields per-slice attention."),
    ("Classifier head", "B × 3 logits", "common",
     "→ normal / near-normal / abnormal  (+ Grad-CAM++ on top-attended slices)."),
]
pipeline_slide("Capped slices  (default)",
               "Bounded bag size → fast, batchable; relies on inter-slice redundancy.", capped)

# ---------------- Slide 3: all slices ----------------
allsl = [
    ("DICOM set", "prepared_brain/<studyID>_<seriesID>/", "io",
     "N slices per set — VARIES (18 … 300+). One label per set."),
    ("index_studies()", "(study_id, [slice paths], label)", "common",
     "labels.csv: study_id,label  (normal / near-normal / abnormal)"),
    ("Slice selection — ALL", "all_slices=True  →  use every slice", "diff",
     "K = N (no cap). Full anatomical coverage, no subsampling."),
    ("Per-slice render → 3 channels", "HU → brain · subdural · bone windows", "common",
     "resize 224×224 + ImageNet norm  →  one slice = 3 × 224 × 224"),
    ("Bag = stack(N slices)", "N × 3 × 224 × 224  (N up to 300+)", "common",
     "Anatomical order (sorted by InstanceNumber)."),
    ("Collate — batch_size = 1", "1 × N × 3 × 224 × 224 · mask all True", "diff",
     "Single bag per step → NO padding waste. grad_accum_steps recovers effective batch."),
    ("MaxViT encoder — CHUNKED + checkpoint", "loop chunks of slice_chunk", "diff",
     "Memory O(chunk), NOT O(N). grad-checkpoint recomputes in backward; AMP autocast halves it."),
    ("Gated-attention pooling", "mask→ -inf, softmax, weighted sum → 1 × D", "common",
     "Collapses N slices → 1 study vector. Same head; attention over ALL real slices."),
    ("Classifier head", "1 × 3 logits", "common",
     "→ normal / near-normal / abnormal  (+ Grad-CAM++ on top-attended slices)."),
]
pipeline_slide("All slices",
               "Every slice used; memory decoupled from N via batch_size=1, chunking, "
               "grad-checkpoint, AMP.", allsl)

# ---------------- Slide 4: comparison ----------------
s = prs.slides.add_slide(BLANK); bg(s)
textbox(s, 0.5, 0.2, 12.4, 0.5, "Capped vs All-slices — where they differ", 26, TITLE, bold=True)

rows = [
    ("Aspect", "Capped (default)", "All slices"),
    ("Slices used", "min(cap, N): 24 train / 32 eval", "all N (no cap)"),
    ("Bag size K", "≤ 32, ~uniform across sets", "= N  (18 – 300+, variable)"),
    ("Batching", "batch_size 8 → pad to maxK", "batch_size 1 (no pad) + grad_accum"),
    ("Encoder pass", "single forward over B·K", "chunked loop + grad-checkpoint"),
    ("Activation memory", "∝ B · cap", "∝ chunk  (independent of N)"),
    ("Precision", "AMP optional", "AMP recommended"),
    ("Speed / epoch", "faster", "slower (recompute + more slices)"),
    ("Coverage", "subsampled (redundancy ok)", "full, no information loss"),
]
left, top, w0, w1, w2 = 0.5, 0.95, 3.1, 4.5, 4.7
rh = 0.46
table = s.shapes.add_table(len(rows), 3, Inches(left), Inches(top),
                           Inches(w0 + w1 + w2), Inches(rh * len(rows))).table
table.columns[0].width = Inches(w0)
table.columns[1].width = Inches(w1)
table.columns[2].width = Inches(w2)
for r, (a, b, c) in enumerate(rows):
    for col, val in enumerate((a, b, c)):
        cell = table.cell(r, col)
        cell.text = val
        para = cell.text_frame.paragraphs[0]
        para.runs[0].font.size = Pt(11.5)
        para.runs[0].font.name = "Calibri"
        if r == 0:
            cell.fill.solid(); cell.fill.fore_color.rgb = TITLE
            para.runs[0].font.color.rgb = WHITE; para.runs[0].font.bold = True
        else:
            cell.fill.solid()
            cell.fill.fore_color.rgb = WHITE if r % 2 else RGBColor(0xEC, 0xF1, 0xF7)
            para.runs[0].font.color.rgb = NOTE
            if col == 2:
                para.runs[0].font.color.rgb = DIFF

# footer: 3-layer handling + flags
y = top + rh * len(rows) + 0.25
box(s, 0.5, y, 12.3, 0.95,
    "Variable slice count is handled at 3 layers",
    "common",
    "(1) selection: cap or all → K   →   (2) collate: pad + mask → rectangular tensor   →   "
    "(3) masked attention: K → 1 fixed study embedding")
textbox(s, 0.5, y + 1.05, 12.3, 0.7,
        "All-slices flags:  --all_slices true  --batch_size 1  --slice_chunk 16  "
        "--grad_checkpoint true  --use_amp true  --grad_accum_steps 8       "
        "(cheapest: --freeze_backbone true → encode under no_grad, train head only)",
        11.5, NOTE, italic=True)

out = "/root/ritikkumar/ct_brain_classifier/docs/ct_brain_flow.pptx"
prs.save(out)
print("saved:", out, "| slides:", len(prs.slides._sldIdLst))
